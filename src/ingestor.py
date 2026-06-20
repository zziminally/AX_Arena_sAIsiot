import re
import pdfplumber
import pandas as pd
from pathlib import Path
from pptx import Presentation
from typing import List, Dict, Any, Optional

from src.config import PROPOSALS_DIR, SAMPLE_DIR, METADATA_CSV, BOILERPLATE

PAGE_NUM_RE    = re.compile(r"^\d{1,2}$")
SECTION_NUM_RE = re.compile(r"^\d{2}\s+")   # "01 " prefix on section divider slides
ACT_RE         = re.compile(r"^ACT\s+0\d")  # "ACT 01" … "ACT 05"


# ── Metadata ─────────────────────────────────────────────────────────────────

def load_metadata() -> Dict[str, Dict]:
    """CSV → {doc_id: metadata_dict}. Pipe-separated fields → lists."""
    df = pd.read_csv(METADATA_CSV)
    list_fields = ["주요 니즈 키워드", "포함 섹션", "강조 가치"]
    for f in list_fields:
        df[f] = df[f].fillna("").apply(
            lambda x: [v.strip() for v in x.split("|") if v.strip()]
        )
    meta: Dict[str, Dict] = {}
    for _, row in df.iterrows():
        doc_id = row["문서 ID"]
        meta[doc_id] = {
            "doc_id":             doc_id,
            "industry":           row["산업군"],
            "client_type":        row["고객사 유형"],
            "proposal_purpose":   row["제안 목적"],
            "project_type":       row["프로젝트 유형"],
            "doc_year":           int(row["문서 연도"]),
            "proposal_stage":     row["제안 단계"],
            "key_needs_keywords": row["주요 니즈 키워드"],
            "included_sections":  row["포함 섹션"],
            "tone_and_manner":    row["톤앤매너"],
            "emphasized_values":  row["강조 가치"],
        }
    return meta


# ── PPTX parsing ─────────────────────────────────────────────────────────────

def _slide_text(slide) -> str:
    """Flatten all paragraph text from a slide, filtering boilerplate."""
    parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t and t not in BOILERPLATE and not PAGE_NUM_RE.match(t):
                parts.append(t)
    return "\n".join(parts)


def _is_divider(text: str) -> bool:
    """
    True for section title cards: total text < 100 chars AND ≤ 3 non-empty lines.
    """
    lines = [l for l in text.splitlines() if l.strip()]
    return len(text) < 100 and len(lines) <= 3


def _extract_section_name(divider_text: str, first_content_text: str) -> str:
    """
    Derive a human-readable section name.
    Priority: first line of the first content slide (always contains the section header).
    Fallback: divider text with numeric prefix stripped.
    This handles both "01 제안 배경..." dividers and "ACT 02" style dividers.
    """
    if first_content_text:
        first_line = first_content_text.splitlines()[0].strip()
        # Content slide first line is the section header — use it directly
        if first_line and len(first_line) <= 50:
            return first_line
    # Fallback: strip "01 " style prefix from divider text
    raw = divider_text.splitlines()[0].strip()
    return SECTION_NUM_RE.sub("", raw) or raw


def parse_pptx(pptx_path: Path) -> List[Dict[str, Any]]:
    """
    PPTX → section-level chunks.
    Slides 0-1 (cover + TOC) → single '표지/목차' chunk.
    Slide 2+: section dividers delimit groups; section name is taken from
    the first content slide (robust across '01 섹션명' and 'ACT 0N' divider styles).
    """
    prs = Presentation(str(pptx_path))
    texts = [_slide_text(s) for s in prs.slides]

    cover_text = "\n".join(t for t in texts[:2] if t)
    chunks: List[Dict] = [
        {"section_name": "표지/목차", "section_order": 0, "text": cover_text}
    ]

    # Collect (divider_text, [content_texts]) groups
    groups: List[tuple] = []  # (divider_text, [content_texts])
    current_divider: Optional[str] = None
    current_body: List[str] = []

    for text in texts[2:]:
        if not text:
            continue
        if _is_divider(text):
            if current_divider is not None:
                groups.append((current_divider, current_body))
            current_divider = text
            current_body = []
        else:
            current_body.append(text)

    if current_divider is not None:
        groups.append((current_divider, current_body))

    for order, (divider_text, body_texts) in enumerate(groups, start=1):
        if not body_texts:
            continue
        section_name = _extract_section_name(divider_text, body_texts[0])
        chunks.append({
            "section_name":  section_name,
            "section_order": order,
            "text":          "\n".join(body_texts),
        })

    return chunks


# ── PDF parsing (회사소개서) ──────────────────────────────────────────────────

_ACT_TITLES = {
    "ACT 01": "광고 시장의 새로운 패러다임",
    "ACT 02": "고객을 움직이는 체험 마케팅",
    "ACT 03": "No.1 광고+테크 컴퍼니 올림플래닛",
    "ACT 04": "독보적인 고객경험 및 레퍼런스",
    "ACT 05": "지금 바로 시작하세요",
}


def parse_company_pdf(pdf_path: Path) -> List[Dict[str, Any]]:
    """
    회사소개서 PDF → ACT 단위 5개 청크.
    각 페이지의 "ACT 0N" 접두어로 소속 ACT를 판별하고,
    같은 ACT의 모든 페이지 텍스트를 하나의 청크로 병합.
    """
    act_buckets: Dict[str, List[str]] = {k: [] for k in _ACT_TITLES}
    current_act: Optional[str] = None

    with pdfplumber.open(str(pdf_path)) as pdf:
        for page in pdf.pages:
            t = (page.extract_text() or "").strip()
            t = re.sub(r"OLIM PLANET.*?Inc\.\s*\d*", "", t).strip()
            if not t:
                continue
            first_line = t.splitlines()[0].strip()
            # 새 ACT 시작 판별: 페이지 첫 줄이 "ACT 0N" 패턴
            if ACT_RE.match(first_line):
                act_key = first_line[:6]
                if act_key in act_buckets:
                    current_act = act_key
            if current_act:
                act_buckets[current_act].append(t)

    chunks: List[Dict] = []
    for order, (act_key, pages) in enumerate(_ACT_TITLES.items()):
        text = "\n\n".join(act_buckets.get(act_key, []))
        if text.strip():
            chunks.append({
                "section_name":  f"{act_key}: {_ACT_TITLES[act_key]}",
                "section_order": order,
                "text":          text,
            })
    return chunks


# ── Chunk builder ─────────────────────────────────────────────────────────────

def _make_chroma_meta(base_meta: Dict, sec: Dict, file_path: str) -> Dict:
    """Flatten metadata for ChromaDB (no lists — primitives only)."""
    m = base_meta
    return {
        "doc_id":             m["doc_id"],
        "section_name":       sec["section_name"],
        "section_order":      sec["section_order"],
        "industry":           m["industry"],
        "client_type":        m["client_type"],
        "proposal_purpose":   m["proposal_purpose"],
        "project_type":       m["project_type"],
        "doc_year":           m["doc_year"],
        "proposal_stage":     m["proposal_stage"],
        "tone_and_manner":    m["tone_and_manner"],
        "key_needs_keywords": "|".join(m.get("key_needs_keywords", [])),
        "included_sections":  "|".join(m.get("included_sections", [])),
        "emphasized_values":  "|".join(m.get("emphasized_values", [])),
        "file_path":          file_path,
    }


def build_chunks(metadata: Dict[str, Dict]) -> List[Dict[str, Any]]:
    """
    1) PPTX files in mock_data/ → proposal chunks
    2) Company brochure PDF in sample/ → company profile chunks
    """
    all_chunks: List[Dict] = []
    seen: set = set()

    # ── Proposal PPTX ────────────────────────────────────────────────────────
    for pptx_path in sorted(PROPOSALS_DIR.glob("*.pptx")):
        doc_id = pptx_path.stem.split("_")[0]
        if doc_id in seen:
            continue
        seen.add(doc_id)

        if doc_id not in metadata:
            print(f"  [skip] {doc_id}: not found in metadata CSV")
            continue

        for sec in parse_pptx(pptx_path):
            if not sec["text"].strip():
                continue
            chunk_id = f"{doc_id}_sec_{sec['section_order']:02d}"
            all_chunks.append({
                "chunk_id": chunk_id,
                "text":     sec["text"],
                "metadata": _make_chroma_meta(metadata[doc_id], sec, str(pptx_path)),
            })

    # ── Company brochure PDF ─────────────────────────────────────────────────
    company_meta = {
        "doc_id":             "COMPANY_PROFILE",
        "industry":           "올림플래닛",
        "client_type":        "올림플래닛 자체",
        "proposal_purpose":   "회사 소개 및 레퍼런스",
        "project_type":       "회사소개",
        "doc_year":           2026,
        "proposal_stage":     "참고자료",
        "tone_and_manner":    "프리미엄·전략적·설득형",
        "key_needs_keywords": ["체험마케팅", "XR", "브랜드캠페인", "올림플래닛"],
        "included_sections":  ["회사소개", "시장분석", "제품소개", "레퍼런스"],
        "emphasized_values":  ["브랜드 신뢰", "기술 이해", "레퍼런스"],
    }
    for pdf_path in SAMPLE_DIR.glob("*.pdf"):
        for sec in parse_company_pdf(pdf_path):
            if not sec["text"].strip():
                continue
            chunk_id = f"COMPANY_PROFILE_sec_{sec['section_order']:02d}"
            all_chunks.append({
                "chunk_id": chunk_id,
                "text":     sec["text"],
                "metadata": _make_chroma_meta(company_meta, sec, str(pdf_path)),
            })

    return all_chunks
