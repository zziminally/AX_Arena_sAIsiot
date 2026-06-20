"""
PPT Assembler: copies top-1 retrieved PPTX, replaces text while preserving
ALL formatting (color, font, size, spacing, bullet style) via XML deepcopy.
"""
import re
import shutil
from copy import deepcopy
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from src.config import OUTPUT_DIR, BOILERPLATE
from src.ingestor import _slide_text, _is_divider, _extract_section_name

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PAGE_NUM_RE = re.compile(r"^\d{1,3}$")
_ACT_RE     = re.compile(r"^ACT\s*(\d+)$", re.IGNORECASE)

# 브랜딩 텍스트 감지 (contains 방식 — full string match 불필요)
_BRANDING_CONTAIN = re.compile(r"OLIMPLANET|©\s*OLIMPLANET", re.IGNORECASE)

# 슬라이드 상단 1400pt 이하 = 헤드라인/섹션 레이블 영역
HEADLINE_ZONE_MAX_EMU = 1_400_000

# 카드 레이아웃 컬럼 그룹핑 허용 폭 (300pt)
CARD_COL_TOLERANCE_PT = 300

# 마침표·쉼표 앞 공백 제거
_PUNCT_RE = re.compile(r"\s+([.。,，!?！？])")


def _clean(text: str) -> str:
    return _PUNCT_RE.sub(r"\1", text).strip()


# ── Overflow 방지 ─────────────────────────────────────────────────────────────

def _enable_autofit(tf) -> None:
    txBody = tf._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is None:
        return
    for tag in [qn("a:normAutofit"), qn("a:spAutoFit"), qn("a:noAutofit")]:
        elem = bodyPr.find(tag)
        if elem is not None:
            bodyPr.remove(elem)
    etree.SubElement(bodyPr, qn("a:normAutofit"))


# ── Format-preserving text replacement ───────────────────────────────────────

def _set_para_text(para_elem, text: str) -> None:
    text = _clean(text)
    runs = para_elem.findall(f"{{{NS}}}r")
    brs  = para_elem.findall(f"{{{NS}}}br")
    for elem in runs + brs:
        para_elem.remove(elem)
    if not runs:
        r = etree.SubElement(para_elem, qn("a:r"))
        t = etree.SubElement(r, qn("a:t"))
        t.text = text
        return
    template_r = deepcopy(runs[0])
    t_elem = template_r.find(f"{{{NS}}}t")
    if t_elem is None:
        t_elem = etree.SubElement(template_r, qn("a:t"))
    t_elem.text = text
    para_elem.append(template_r)


def _set_tf_single(tf, text: str) -> None:
    txBody = tf._txBody
    paras = txBody.findall(f"{{{NS}}}p")
    if not paras:
        return
    for p in paras[1:]:
        txBody.remove(p)
    _set_para_text(paras[0], text)
    _enable_autofit(tf)


def _set_tf_lines(tf, lines: List[str]) -> None:
    if not lines:
        return
    txBody = tf._txBody
    paras = txBody.findall(f"{{{NS}}}p")
    if not paras:
        return
    template_p = next((p for p in paras if p.findall(f"{{{NS}}}r")), paras[0])
    for p in paras:
        txBody.remove(p)
    for line in lines:
        new_p = deepcopy(template_p)
        _set_para_text(new_p, line)
        txBody.append(new_p)
    _enable_autofit(tf)


def _set_tf_with_breaks_v2(tf, text: str) -> None:
    """\\n → soft line break within one paragraph, preserving run format."""
    parts = [_clean(p) for p in re.split(r"\\n|\n", text) if p.strip()]
    if not parts:
        return
    txBody = tf._txBody
    paras = txBody.findall(f"{{{NS}}}p")
    if not paras:
        return
    para = paras[0]
    for p in paras[1:]:
        txBody.remove(p)
    runs = para.findall(f"{{{NS}}}r")
    template_r = deepcopy(runs[0]) if runs else None
    for elem in para.findall(f"{{{NS}}}r") + para.findall(f"{{{NS}}}br"):
        para.remove(elem)
    for i, part in enumerate(parts):
        if i > 0:
            br = etree.SubElement(para, qn("a:br"))
            if template_r is not None:
                rpr = template_r.find(f"{{{NS}}}rPr")
                if rpr is not None:
                    br.insert(0, deepcopy(rpr))
        if template_r is not None:
            new_r = deepcopy(template_r)
            t_elem = new_r.find(f"{{{NS}}}t")
            if t_elem is None:
                t_elem = etree.SubElement(new_r, qn("a:t"))
            t_elem.text = part
            para.append(new_r)
        else:
            r = etree.SubElement(para, qn("a:r"))
            t = etree.SubElement(r, qn("a:t"))
            t.text = part
    _enable_autofit(tf)


# ── Branding cleanup ──────────────────────────────────────────────────────────

def _clear_text_shape(shape) -> None:
    """Run 텍스트를 모두 빈 문자열로 초기화."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""


def _clear_slide_branding(slide, slide_height: int) -> None:
    """
    모든 슬라이드에서 OLIMPLANET 워터마크, © 푸터, 하단 페이지 번호를 제거.
    섹션 레이블(빨간 텍스트)·카드 번호 등은 건드리지 않음.
    """
    bottom_threshold = int(slide_height * 0.85)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if not t:
            continue
        if _BRANDING_CONTAIN.search(t):
            _clear_text_shape(shape)
        elif PAGE_NUM_RE.match(t) and shape.top > bottom_threshold:
            # 하단 페이지 번호만 제거 (카드 내부 번호는 top이 낮으므로 제외)
            _clear_text_shape(shape)


# ── Shape classification ──────────────────────────────────────────────────────

def _text_shapes_by_area(slide) -> List[Tuple[int, Any]]:
    result = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if not t or t in BOILERPLATE or PAGE_NUM_RE.match(t):
            continue
        result.append((shape.width * shape.height, shape))
    return sorted(result, key=lambda x: -x[0])


# ── Cover slide ───────────────────────────────────────────────────────────────

def _replace_cover(slide, cover: Dict, slide_height: int) -> None:
    """
    커버 슬라이드: 브랜딩 제거 → title/subtitle 교체.
    _clear_slide_branding으로 PROPOSAL, OLIMPLANET, © 전부 제거.
    """
    _clear_slide_branding(slide, slide_height)

    # PROPOSAL 같은 별도 텍스트 shape도 직접 지움
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if t.upper() == "PROPOSAL":
            _clear_text_shape(shape)

    shapes = _text_shapes_by_area(slide)
    if not shapes:
        return
    title    = cover.get("title", "")
    subtitle = cover.get("subtitle", "")
    if len(shapes) >= 2:
        _set_tf_single(shapes[0][1].text_frame, title)
        _set_tf_with_breaks_v2(shapes[1][1].text_frame, subtitle)
    elif shapes:
        _set_tf_single(shapes[0][1].text_frame, title)
    for _, shape in shapes[2:]:
        _set_tf_single(shape.text_frame, "")


# ── Section divider slide ─────────────────────────────────────────────────────

def _replace_divider_slide(slide, act_tagline: str) -> None:
    """
    섹션 구분 슬라이드:
    - 가장 큰 shape → act_tagline
    - 'ACT NN' shape → 'NN' (ACT 접두어 제거)
    """
    shapes = _text_shapes_by_area(slide)
    if not shapes:
        return
    _set_tf_single(shapes[0][1].text_frame, act_tagline)
    for _, shape in shapes[1:]:
        t = shape.text_frame.text.strip()
        m = _ACT_RE.match(t)
        if m:
            _set_tf_single(shape.text_frame, m.group(1))
            break


# ── Content slide ─────────────────────────────────────────────────────────────

def _replace_content_slide(slide, sec_draft: Dict, slide_height: int = 6858000) -> None:
    """
    Y좌표 기반 shape 분류:

    hl_zone   (top < HEADLINE_ZONE_MAX_EMU):  가장 큰 shape → headline
    body_zone (top ≥ HEADLINE_ZONE_MAX_EMU):
      - 1개     → 표준 body (subtitle + bullets)
      - 복수 / 카드 레이아웃:
          300pt 컬럼 허용폭으로 묶어 title+description 쌍 감지
          → "항목명: 설명" 분리해 title shape=항목명, desc shape=설명
    """
    shapes = _text_shapes_by_area(slide)
    if not shapes:
        return

    headline = sec_draft.get("headline", "")
    subtitle = sec_draft.get("subtitle", "")
    bullets  = sec_draft.get("bullets", [])

    body_lines: List[str] = []
    if subtitle:
        body_lines.extend([p.strip() for p in re.split(r"\\n|\n", subtitle) if p.strip()])
    body_lines.extend(bullets)

    # 단일 shape → 전부 삽입
    if len(shapes) == 1:
        all_lines = ([headline] if headline else []) + body_lines
        _set_tf_lines(shapes[0][1].text_frame, all_lines)
        return

    hl_zone   = sorted([(a, s) for a, s in shapes if s.top <  HEADLINE_ZONE_MAX_EMU], key=lambda x: -x[0])
    body_zone = sorted([(a, s) for a, s in shapes if s.top >= HEADLINE_ZONE_MAX_EMU], key=lambda x: -x[0])

    # 헤드라인
    headline_shape = hl_zone[0][1] if hl_zone else shapes[0][1]
    _set_tf_single(headline_shape.text_frame, headline)

    if not body_zone:
        return

    if len(body_zone) == 1:
        _set_tf_lines(body_zone[0][1].text_frame, body_lines)
        return

    # 복수 body shapes → 카드 레이아웃 처리
    # round/floor 기반 그리드 대신 인접 거리 클러스터링 사용:
    # 같은 카드 내 title·description shape의 left 차이 ≈ 60pt,
    # 카드 간 gap ≈ 2500pt → 150pt 임계값으로 정확히 구분
    GAP_EMU = CARD_COL_TOLERANCE_PT * 914

    sorted_shapes = sorted(body_zone, key=lambda x: x[1].left)
    col_groups: List[List[Tuple[int, Any]]] = []
    cur: List[Tuple[int, Any]] = [sorted_shapes[0]]
    for item in sorted_shapes[1:]:
        if item[1].left - cur[-1][1].left <= GAP_EMU:
            cur.append(item)
        else:
            col_groups.append(cur)
            cur = [item]
    col_groups.append(cur)

    for i, col_shapes in enumerate(col_groups):
        if i >= len(bullets):
            break
        bullet = bullets[i]

        # "항목명: 설명" 분리
        if ": " in bullet:
            card_title, card_body = bullet.split(": ", 1)
        else:
            card_title, card_body = bullet, bullet

        # 컬럼 내 면적 오름차순: 작은 shape=제목, 큰 shape=설명
        col_sorted = sorted(col_shapes, key=lambda x: x[0])
        if len(col_sorted) >= 2:
            _set_tf_single(col_sorted[0][1].text_frame, card_title)
            _set_tf_single(col_sorted[-1][1].text_frame, card_body)
        else:
            _set_tf_single(col_sorted[0][1].text_frame, bullet)


# ── Section-to-slide mapping ──────────────────────────────────────────────────

def _normalize(name: str) -> str:
    return re.sub(r"[\s/·\-또는|]", "", name).lower()


def _map_sections_to_slides(prs: Presentation) -> Dict[str, Dict]:
    texts = [_slide_text(s) for s in prs.slides]
    section_map: Dict[str, Dict] = {}
    current_section: Optional[str] = None
    for i, text in enumerate(texts):
        if i < 2:
            continue
        if not text:
            continue
        if _is_divider(text):
            next_text = texts[i + 1] if i + 1 < len(texts) else ""
            current_section = _extract_section_name(text, next_text)
            section_map[current_section] = {"divider": i, "content": []}
        elif current_section is not None:
            section_map[current_section]["content"].append(i)
    return section_map


def _find_section_info(sec_name: str, section_map: Dict[str, Dict]) -> Dict:
    _empty = {"divider": None, "content": []}
    if sec_name in section_map:
        return section_map[sec_name]
    norm_query = _normalize(sec_name)
    for mapped, info in section_map.items():
        if _normalize(mapped) == norm_query:
            return info
    best_info, best_overlap = _empty, 0
    for mapped, info in section_map.items():
        norm_mapped = _normalize(mapped)
        overlap = sum(1 for k in range(len(norm_query) - 2) if norm_query[k:k+3] in norm_mapped)
        if overlap > best_overlap:
            best_overlap, best_info = overlap, info
    return best_info if best_overlap >= 1 else _empty


# ── Main assembly ─────────────────────────────────────────────────────────────

def assemble(draft: Dict[str, Any], template_pptx_path: str) -> Path:
    OUTPUT_DIR.mkdir(exist_ok=True)
    title_slug = re.sub(r"[^\w가-힣]", "_", draft["cover"].get("title", "proposal"))[:40]
    output_path = OUTPUT_DIR / f"draft_{title_slug}.pptx"

    shutil.copy(template_pptx_path, output_path)
    prs = Presentation(str(output_path))
    slide_h = prs.slide_height

    # 1. Cover (브랜딩 제거는 커버에만 적용)
    _replace_cover(prs.slides[0], draft["cover"], slide_h)

    # 3. Section map
    section_map = _map_sections_to_slides(prs)
    slide2_lookup: Dict[str, Dict] = {
        s.get("section_name", ""): s for s in draft.get("sections_slide2", [])
    }

    # 4. Replace per section
    for sec_draft in draft.get("sections", []):
        sec_name = sec_draft.get("section_name", "")
        info = _find_section_info(sec_name, section_map)

        if info["divider"] is not None and sec_draft.get("act_tagline"):
            _replace_divider_slide(prs.slides[info["divider"]], sec_draft["act_tagline"])

        content = info["content"]
        if not content:
            continue

        _replace_content_slide(prs.slides[content[0]], sec_draft, slide_h)

        if len(content) >= 2:
            s2 = slide2_lookup.get(sec_name)
            if s2:
                _replace_content_slide(prs.slides[content[1]], s2, slide_h)

    prs.save(str(output_path))
    return output_path
